#!/usr/bin/env python3
"""Z-MAX PPT — 智蜂模板风格 · DDS全站数据 · 白底蓝强调 · 专业商务风"""
import sqlite3, os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

TEMPLATE = '/www/wwwroot/datadrive.world/uploads/data/智蜂具身机器人产品规划和宣传页_v1.1_0422.pptx'
OUT = '/www/wwwroot/datadrive.world/Z700-立项申请书.pptx'

# ====== 智蜂模板设计系统（从实际模板分析提取）======
BLUE   = RGBColor(0x00, 0x66, 0xCC)    # 强调蓝（模板主色）
DARK   = RGBColor(0x00, 0x00, 0x00)    # 主标题黑
GRAY1  = RGBColor(0x1F, 0x29, 0x37)    # 副标题深灰
GRAY2  = RGBColor(0x33, 0x33, 0x33)    # 卡片标题
GRAY3  = RGBColor(0x6B, 0x72, 0x80)    # 描述灰
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)    # 白
ORANGE = RGBColor(0xF9, 0x73, 0x16)    # 关键数字强调
LIGHT_BORDER = RGBColor(0xE5, 0xE7, 0xEB)  # 浅边框
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)  # 卡片底色

# 模板画布: 26.7" x 15.0" (超宽屏 16:9)
W = 26.7; H = 15.0

def load():
    db = sqlite3.connect('/www/wwwroot/datadrive.world/dds.db')
    db.row_factory = sqlite3.Row
    d = {}
    d['c'] = {r['key']: r['value'] for r in db.execute('SELECT key,value FROM company')}
    d['kpi'] = {r['id']: dict(r) for r in db.execute('SELECT * FROM kpi')}
    d['robots'] = [dict(r) for r in db.execute('SELECT * FROM robots ORDER BY id')]
    d['roadmap'] = [dict(r) for r in db.execute('SELECT * FROM roadmap ORDER BY version')]
    d['zones'] = [dict(r) for r in db.execute('SELECT * FROM factory_zones ORDER BY id')]
    d['pipe'] = [dict(r) for r in db.execute('SELECT * FROM pipeline ORDER BY step')]
    d['st'] = db.execute('SELECT COUNT(*) FROM atomic_skills').fetchone()[0]
    d['sc'] = {}
    for r in db.execute('SELECT category, COUNT(*) as cnt FROM atomic_skills GROUP BY category ORDER BY cnt DESC'):
        d['sc'][r['category']] = r['cnt']
    db.close()
    return d

D = load(); c = D['c']

# ====== 复制模板，保留所有原始页作为附录 ======
shutil.copy(TEMPLATE, OUT)
prs = Presentation(OUT)

# 品牌替换（模板原有品牌名 → 智蜂创元）
for sl in list(prs.slides):
    for sh in list(sl.shapes):
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for run in p.runs:
                    if 'ReinBee' in run.text:
                        run.text = run.text.replace('ReinBee', c.get('name', '智蜂创元'))

CONTENT_LAYOUT = prs.slide_layouts[1]  # 继承模板背景图

def ns(idx=None):
    """在指定位置插入新slide（使用模板布局，继承背景图）。
       不传idx则追加到末尾。"""
    if idx is not None:
        # 插入到指定位置：先追加临时slide，再移动XML顺序
        sl = prs.slides.add_slide(CONTENT_LAYOUT)
        sldIdLst = list(prs.slides._sldIdLst)
        # 把最后一个（刚加的）移到idx位置
        last = sldIdLst.pop()
        sldIdLst.insert(idx, last)
        # 重建 _sldIdLst
        prs.slides._sldIdLst.clear()
        for elem in sldIdLst:
            prs.slides._sldIdLst.append(elem)
        return sl
    else:
        return prs.slides.add_slide(CONTENT_LAYOUT)

# ====== UI 组件函数 ======

def add_title(slide, text, left=1.2, top=0.5, width=24.5, sz=Pt(44)):
    """主标题 — 黑色粗体（匹配模板 48pt 风格）"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.bold = True; p.font.color.rgb = DARK
    return tb

def add_sub(slide, text, top=1.4, sz=Pt(24)):
    """副标题 — 深灰色"""
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(top), Inches(24.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.color.rgb = GRAY1
    return tb

def blue_bar(slide, left, top, width, height=Pt(4)):
    """蓝色装饰条"""
    bar = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), height)
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    return bar

def kpi_card(slide, left, top, w, h, value, unit, label, icon=''):
    """KPI 卡片 — 白底+蓝顶边，匹配模板风格"""
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD_BG
    sh.line.color.rgb = LIGHT_BORDER; sh.line.width = Pt(0.8)
    sh.shadow.inherit = False
    
    blue_bar(slide, left, top, w)
    
    # 数值
    tb = slide.shapes.add_textbox(Inches(left), Inches(top+0.3), Inches(w), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = f'{icon} {value}' if icon else value
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = BLUE
    
    # 单位
    tb2 = slide.shapes.add_textbox(Inches(left), Inches(top+2.0), Inches(w), Inches(0.4))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    p2.text = unit; p2.font.size = Pt(16); p2.font.color.rgb = GRAY3
    
    # 标签
    tb3 = slide.shapes.add_textbox(Inches(left), Inches(top+2.4), Inches(w), Inches(0.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    p3.text = label; p3.font.size = Pt(14); p3.font.color.rgb = GRAY2

def info_card(slide, left, top, w, h, title_text, desc_text, title_color=BLUE):
    """信息卡片 — 白底+蓝左边条，匹配模板feature card风格"""
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LIGHT_BORDER; sh.line.width = Pt(0.8)
    
    blue_bar(slide, left, top, 0.08, h)  # 左边蓝条
    
    # 标题
    tb = slide.shapes.add_textbox(Inches(left+0.35), Inches(top+0.2), Inches(w-0.5), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title_text
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = title_color
    
    # 描述
    tb2 = slide.shapes.add_textbox(Inches(left+0.35), Inches(top+0.8), Inches(w-0.55), Inches(h-1.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = desc_text
    p2.font.size = Pt(15); p2.font.color.rgb = GRAY3; p2.line_spacing = Pt(24)

# ============================================================
# 正文 8 页 — 插在模板封面(0)和目录(1)之后
# ============================================================

# ------ S1 封面（index 2，模板封面和目录之后第一页） ------
sl = ns(idx=2)
add_title(sl, 'Z700 光模块精密制造', top=4.5, sz=Pt(52))
add_sub(sl, '具身智能机器人系统 · 立项申请书', top=6.5, sz=Pt(32))
add_sub(sl, f'{c.get("name","智蜂创元")} · datadrive.world', top=8.0, sz=Pt(24))
add_sub(sl, '掌握智能，蜂动未来', top=9.0, sz=Pt(18))
blue_bar(sl, 5.0, 10.5, 16.7)

# ------ S2 KPI 指标（index 3） ------
sl = ns(idx=3)
add_title(sl, '一、核心性能指标', sz=Pt(40))
add_sub(sl, 'FW Loading · 微米级定位 + 柔性力控 · 精密操作试金石', sz=Pt(20))

kpi_items = [('precision', '🎯'), ('yield_rate', '✅'), ('cycle_time', '⏱️'), ('force_bw', '⚡')]
cw, ch, cgap = 5.8, 4.5, 0.5
cx = (W - (cw*4 + cgap*3)) / 2
for i, (k, icon) in enumerate(kpi_items):
    v = D['kpi'].get(k, {})
    kpi_card(sl, cx + i*(cw+cgap), 2.5, cw, ch,
             v.get('value', ''), v.get('unit', ''), v.get('label', k), icon)

tb = sl.shapes.add_textbox(Inches(1.5), Inches(7.5), Inches(23.7), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
p.text = f'数据来源: DDS 全局数据空间 · {D["st"]} 项原子技能 · 63 页面全站同步'
p.font.size = Pt(14); p.font.color.rgb = GRAY3

# ------ S3 产品系列（index 4） ------
sl = ns(idx=4)
add_title(sl, '二、Z700 产品系列', sz=Pt(40))
add_sub(sl, 'L2 → L4 软件定义进化 · 5 款机器人覆盖全场景', sz=Pt(20))

bots = D['robots']
for i, bot in enumerate(bots):
    y = 2.2 + i * 2.3
    tc = ORANGE if bot.get('level') == 'L4' else BLUE
    info_card(sl, 1.5, y, 23.7, 2.0,
              f'{bot["icon"]}  {bot["name"]}  [{bot["level_label"]}]',
              bot.get('desc', ''), tc)

# ------ S4 创元 XWorld（index 5） ------
sl = ns(idx=5)
add_title(sl, '三、创元 XWorld · 全技能智能体平台', sz=Pt(40))

xworld = [
    ('🧠  XAgent 大脑',
     '全技能智能体 · 统一指挥调度 · 安全高效自主交互\n视觉语言动作一体化模型 · 端到端感知决策控制'),
    ('📊  XData 数据',
     '全模态数据贯通 · 人机料法环全要素覆盖\nDDS 全局数据空间 · 实时同步 · 数据基石'),
    ('🤖  XRobot 机器人',
     '垂域模型 + 强健肢体 · 感驱控一体\n精密力控 >1kHz · 柔性操作 · 精准执行'),
]
for i, (t, d) in enumerate(xworld):
    info_card(sl, 1.5, 1.8 + i * 3.3, 23.7, 2.9, t, d, BLUE)

tb = sl.shapes.add_textbox(Inches(1.5), Inches(12.5), Inches(23.7), Inches(0.8))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
p.text = f'{D["st"]} 项原子技能 · {len(D["sc"])} 大类 · DDS 全局数据驱动'
p.font.size = Pt(16); p.font.color.rgb = GRAY2

# ------ S5 数据流水线（index 6） ------
sl = ns(idx=6)
add_title(sl, '四、端到端数据流水线', sz=Pt(40))
add_sub(sl, '感知 → 推理 → 执行 → 回流 · 四步闭环', sz=Pt(20))

pipe = D['pipe']
pw = 5.6
for i, p in enumerate(pipe):
    x = 1.8 + i * (pw + 0.5)
    info_card(sl, x, 2.5, pw, 6.5,
              f'{p.get("icon","")}  {p.get("name","")}',
              p.get('desc', ''), ORANGE if i == 0 else BLUE)

# 箭头
for i in range(len(pipe)-1):
    ax = 1.8 + (i+1)*(pw+0.5) - 0.6
    arr = sl.shapes.add_shape(33, Inches(ax), Inches(5.6), Inches(0.5), Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb = BLUE; arr.line.fill.background()

# ------ S6 工厂四区（index 7） ------
sl = ns(idx=7)
add_title(sl, '五、目标工厂 · 800G DR8 光模块产线', sz=Pt(40))

total_stations = sum(z.get('count', 0) for z in D['zones'])
zone_kpi_data = [
    (str(len(D['zones'])), '产线区域', '四区全流程覆盖'),
    (str(total_stations), '工位', '全流程工站'),
    (str(D['st']), '原子技能', f'{len(D["sc"])} 大类'),
    ('10', '台', '一期机器人部署'),
]
for i, (val, unit, label) in enumerate(zone_kpi_data):
    kpi_card(sl, 1.5 + i*6.2, 1.8, 5.8, 3.0, val, unit, label)

for i, z in enumerate(D['zones']):
    col = i % 2; row = i // 2
    x = 1.5 + col * 12.5; y = 5.5 + row * 2.6
    info_card(sl, x, y, 12.0, 2.2,
              z['name'],
              f'{z.get("count", "")} 站 · {z.get("stations", "")}')

# ------ S7 路线图（index 8） ------
sl = ns(idx=8)
add_title(sl, '六、开发路线图 · 2026 → 2028', sz=Pt(40))

roadmap = D['roadmap']
rw = 5.8
for i, r in enumerate(roadmap):
    x = 1.5 + i * (rw + 0.5)
    info_card(sl, x, 2.0, rw, 6.0,
              r['name'], f'时间: {r.get("timeline","")}\n{r.get("desc","")}',
              ORANGE if i == 0 else BLUE)

# 时间轴
blue_bar(sl, 1.5, 8.8, 24.0)
for i in range(len(roadmap)):
    dx = 4.4 + i * (rw + 0.5)
    dot = sl.shapes.add_shape(9, Inches(dx), Inches(8.6), Inches(0.45), Inches(0.45))
    dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE if i == 0 else BLUE
    dot.line.fill.background()

# ------ S8 原子技能（index 9） ------
sl = ns(idx=9)
add_title(sl, f'七、原子技能库 · {D["st"]} 项 · {len(D["sc"])} 大类', sz=Pt(40))

cats = list(D['sc'].items())
cols = 4; cw2 = 5.6; ch2 = 2.3
for i, (cat, cnt) in enumerate(cats):
    col = i % cols; row = i // cols
    x = 1.5 + col * (cw2 + 0.55); y = 1.8 + row * (ch2 + 0.4)
    if y > 12.5: break
    
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(cw2), Inches(ch2))
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = LIGHT_BORDER; sh.line.width = Pt(0.8)
    blue_bar(sl, x, y, cw2)
    
    tb = sl.shapes.add_textbox(Inches(x+0.3), Inches(y+0.35), Inches(cw2-0.6), Inches(0.55))
    p = tb.text_frame.paragraphs[0]; p.text = cat
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = GRAY2
    
    tb2 = sl.shapes.add_textbox(Inches(x+0.3), Inches(y+1.0), Inches(cw2-0.6), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    p2.text = str(cnt); p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = BLUE
    
    tb3 = sl.shapes.add_textbox(Inches(x+0.3), Inches(y+1.7), Inches(cw2-0.6), Inches(0.4))
    p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    p3.text = '条原子技能'; p3.font.size = Pt(12); p3.font.color.rgb = GRAY3

# ------ S9 团队与总结（index 10） ------
sl = ns(idx=10)
add_title(sl, '八、项目团队与立项总结', sz=Pt(40))

team = [
    ('🛠️  xspace 总工', 'Sys架构 · GUI引擎 · Orin部署', '18项 · 88%'),
    ('💻  web PM/前端', 'Sys2训练 · Web全站 · 仿真', '18项 · 94%'),
    ('🔬  小芳 硬件', 'Orin采集 · MAC转发 · 硬件测试', '14项 · 78%'),
]
for i, (name, role, progress) in enumerate(team):
    x = 1.5 + i * 8.3
    info_card(sl, x, 1.8, 7.8, 3.8, name,
              f'{role}\n完成率: {progress}', ORANGE if i == 0 else BLUE)

info_card(sl, 1.5, 6.2, 23.7, 3.8, '📋  立项总结',
          '技术领先: 视触觉混合动作模型 · 端到端感知决策控制 · >1kHz 力控闭环\n'
          '市场刚需: AI算力驱动光模块扩产 · 精密制造人工替代\n'
          '团队完备: 三人全栈 · 50项 88% 完成 · 零阻塞 · kanban 驱动')

# 数据资产标注
tb = sl.shapes.add_textbox(Inches(1.5), Inches(10.5), Inches(23.7), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
p.text = f'数据资产: DDS 全局数据空间 · 63 页面 · {D["st"]} 项原子技能 · {len(D["sc"])} 大类'
p.font.size = Pt(14); p.font.color.rgb = GRAY2

# ------ S10 谢谢（index 11） ------
sl = ns(idx=11)
add_title(sl, '谢  谢', top=5.0, sz=Pt(56))
add_sub(sl, f'{c.get("name","智蜂创元")} · datadrive.world', top=7.2, sz=Pt(28))
add_sub(sl, '掌握智能，蜂动未来', top=8.2, sz=Pt(22))
blue_bar(sl, 5.0, 10.0, 16.7)

# ====== 保存 ======
prs.save(OUT)
sl_count = len(list(prs.slides))
sz_mb = os.path.getsize(OUT) / 1024 / 1024
print(f'✅ {OUT} ({sz_mb:.1f}MB) · {sl_count} slides · {D["st"]} skills')
print(f'   10 页正文 + {sl_count-10} 页模板附录 · 白底蓝强调 · 字号匹配智蜂设计系统')
