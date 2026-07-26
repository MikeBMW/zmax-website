#!/usr/bin/env python3
"""Generate Z700 analysis PPTX - clean white theme"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C = RGBColor(0x00, 0xAA, 0x88)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x22, 0x66, 0xCC)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Top accent bar
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    shape.fill.solid(); shape.fill.fore_color.rgb = C; shape.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(4))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Z700 轮式双臂 · 精品分析报告"
    p.font.size = Pt(42); p.font.bold = True; p.font.color.rgb = C; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "光模块精密制造 · 具身智能旗舰 · 对比它石智航 TARS"
    p2.font.size = Pt(18); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
    p3 = tf.add_paragraph(); p3.text = "智蜂创元 (苏州)  ×  它石智航 TARS (上海)  ·  2026"
    p3.font.size = Pt(14); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(30)
    # Bottom bar
    shape2 = slide.shapes.add_shape(1, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = C; shape2.line.fill.background()

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Top bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    shape.fill.solid(); shape.fill.fore_color.rgb = C; shape.line.fill.background()
    # Title
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
    tf = tx.text_frame; p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = DARK
    # Green underline
    shape2 = slide.shapes.add_shape(1, Inches(0.6), Inches(0.95), Inches(2), Inches(0.04))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = C; shape2.line.fill.background()
    # Bullets
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    tf2 = tx2.text_frame; tf2.word_wrap = True
    for i, line in enumerate(bullets):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(15)
        p2.font.color.rgb = GRAY if not line.startswith('•') else DARK
        p2.space_after = Pt(6)
        p2.space_before = Pt(2)
        if '→' in line or '✅' in line or '🔶' in line:
            p2.font.size = Pt(14)

def add_stats_slide(title, stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    shape.fill.solid(); shape.fill.fore_color.rgb = C; shape.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
    tf = tx.text_frame; p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = DARK
    shape2 = slide.shapes.add_shape(1, Inches(0.6), Inches(0.95), Inches(2), Inches(0.04))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = C; shape2.line.fill.background()
    # Stats as cards
    cols = min(len(stats), 6)
    card_w = Inches(1.8); gap = Inches(0.15)
    start_x = (prs.slide_width - (card_w + gap) * cols + gap) / 2
    for i, (val, unit, label) in enumerate(stats):
        col = i % cols; row = i // cols
        x = start_x + col * (card_w + gap)
        y = Inches(1.5) + row * Inches(2.5)
        # Card bg
        card = slide.shapes.add_shape(1, x, y, card_w, Inches(2.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
        card.line.color.rgb = RGBColor(0xDD, 0xEE, 0xEE)
        # Value
        txv = slide.shapes.add_textbox(x, y + Inches(0.3), card_w, Inches(0.8))
        tfv = txv.text_frame; pv = tfv.paragraphs[0]
        pv.text = val; pv.font.size = Pt(32); pv.font.bold = True; pv.font.color.rgb = C; pv.alignment = PP_ALIGN.CENTER
        # Unit
        txu = slide.shapes.add_textbox(x, y + Inches(1.0), card_w, Inches(0.4))
        tfu = txu.text_frame; pu = tfu.paragraphs[0]
        pu.text = unit; pu.font.size = Pt(12); pu.font.color.rgb = GRAY; pu.alignment = PP_ALIGN.CENTER
        # Label
        txl = slide.shapes.add_textbox(x, y + Inches(1.35), card_w, Inches(0.5))
        tfl = txl.text_frame; pl = tfl.paragraphs[0]
        pl.text = label; pl.font.size = Pt(10); pl.font.color.rgb = GRAY; pl.alignment = PP_ALIGN.CENTER

# ═══ Slide 1: Title ═══
add_title_slide()

# ═══ Slide 2: KPI ═══
add_stats_slide("核心性能指标", [
    ("±0.02", "mm", "重复定位精度"),
    (">99", "%", "关键工序良率"),
    (">10", "kHz", "力控闭环带宽"),
    ("<15", "s", "单次插拔节拍"),
    ("≤0.1", "N", "接触力控制"),
    (">50", "kg", "双臂负载"),
])

# ═══ Slide 3: Product ═══
add_content_slide("Z700 轮式双臂 · 产品架构", [
    "L4 旗舰级 · 全自主执行 + 自适应 + 自恢复",
    "",
    "🔵 感知系统 — 2D/3D视觉 · 六维力/力矩 · 触觉阵列 · 结构光 · 激光雷达 · IMU",
    "    ↳ 原子技能 P001-P015 感知定位 (15项)",
    "",
    "🟢 操作执行 — 双臂6轴×2协同 · 精密夹爪/真空 · 力位混合 · 柔顺装配",
    "    ↳ 原子技能 A001-A015 操作执行 (15项)",
    "",
    "🟡 移动导航 — 轮式AMR底盘 · SLAM · 多工位自由调度 · 精准对接",
    "    ↳ 原子技能 M001-M007 移动导航 (7项)",
    "",
    "🟣 力控安全 — 碰撞检测 · 10kHz力控闭环 · 安全互锁 · 急停",
    "    ↳ 原子技能 S001-S010 安全集成 (10项)",
])

# ═══ Slide 4: Scenarios ═══
add_content_slide("五大应用场景", [
    "✅ Phase 1 (2026.07-11)",
    "  1. FW固件烧录  — 扫码→烧录→校验→追溯                    KPI: 节拍 ≤15s",
    "  2. 精密上下料  — 识别→取件→放置→满空交换                  KPI: 零损伤",
    "  3. 老化箱插拔  — 6D定位→力控插入→锁止→确认              KPI: 力 ≤0.1N",
    "",
    "🔶 Phase 2 (2026.12-2027.04)",
    "  4. 热海柜操作  — 开柜→取放→关柜→循环监控                KPI: ±0.05mm",
    "  5. ATS自动测试 — 装夹→插拔→测试→分Bin→追溯             KPI: ≥99%良率",
    "",
    "覆盖光模块产线: COC基板 → OE光引擎 → MOD模块组装 → 全流程"
])

# ═══ Slide 5: TARS ═══
add_content_slide("它石智航 TARS 对标分析", [
    "智蜂创元(苏州) × 它石智航 TARS(上海)  ·  一期10台(5自研Z700+5台TARS)  ·  1000万预算  ·  白盒交付",
    "",
    "对比维度              │  Z700 智蜂自研          │  TARS 它石智航          │  结论",
    "──────────────────────┼─────────────────────────┼─────────────────────────┼──────────────",
    "力控精度              │  ≤0.1N / 10kHz闭环      │  ≤0.5N / 1kHz典型       │  Z700领先 5×",
    "重复定位精度          │  ±0.02mm                │  ±0.05mm                │  Z700领先 2.5×",
    "插拔节拍              │  <15s/次                │  20-30s/次              │  Z700快 1.5-2×",
    "自主导航              │  AMR多工位自由移动      │  需轨道/固定安装        │  Z700独有优势",
    "双臂协同              │  原生双臂协同装配        │  单臂为主，需定制       │  Z700架构优势",
    "软件生态              │  89原子技能 · 全栈自研  │  ROS2通用框架           │  专用 vs 通用",
    "交付模式              │  白盒 · 全栈技术转让    │  黑盒 · SDK接口         │  Z700战略优势",
    "部署周期              │  2-4周/工站             │  4-8周/工站             │  Z700更快",
])

# ═══ Slide 6: SWOT ═══
add_content_slide("SWOT 分析", [
    "💪 优势 (Strengths)",
    "  • 光模块行业深度定制，89项原子技能覆盖全产线",
    "  • 力控闭环 >10kHz，精密操作能力行业领先",
    "  • 轮式AMR + 双臂架构，多工位自由调度",
    "  • 白盒交付，客户获全栈能力 · 苏州产业集群优势",
    "",
    "⚠️ 劣势 (Weaknesses)",
    "  • 品牌知名度低于它石/珞石等老牌厂商",
    "  • 量产成熟度待 Phase1 验证 · 供应链外部依赖",
    "",
    "🚀 机会 (Opportunities)",
    "  • 光模块 800G/1.6T 升级 + AI算力驱动产能扩张",
    "  • 机器换人政策红利 · 白盒模式可复制半导体/汽车电子",
    "",
    "🔻 威胁 (Threats)",
    "  • 它石可能自研光模块方案 · 成熟厂商向下整合 · 价格战风险",
])

# ═══ Slide 7: Skills ═══
add_content_slide("89项原子技能 × 光模块全产线覆盖", [
    "感知定位 15项  P001-P015   对象识别 · 位姿估计 · 缺陷检测 · 视觉伺服 · 条码/OCR · 3D点云",
    "操作执行 15项  A001-A015   精密夹取 · 力控插入 · 柔顺装配 · 锁付 · 点胶 · 焊接 · 贴装",
    "装配工艺 15项  P056-P070   COC共晶 · WB打线 · UV固化 · 耦合对准 · 分板 · 绑定 · 测试",
    "质量检测 10项  Q001-Q010   目检 · AOI · 3D测量 · GR&R · 缺陷分类 · 追溯报告",
    "安全集成 10项  S001-S010   碰撞检测 · 力控闭环 · 急停 · 区域监控 · 安全互锁",
    "载具物流  9项  H001-H009   载具识别 · 工站接驳 · 满空交换 · 槽位追溯 · 异常恢复",
    "学习泛化  8项  L001-L008   技能迁移 · 自校准 · 少样本学习 · 多模态指令理解",
    "移动导航  7项  M001-M007   SLAM · 路径规划 · 动态避障 · 工站精准对接",
    "",
    "总计: 89项  |  8大类别  |  覆盖COC + OE + MOD + WH 全产线工艺",
])

# ═══ Slide 8: Roadmap ═══
add_content_slide("交付路线图", [
    "Phase 1  (2026.07 — 2026.11)  首批样机 + POC验证",
    "  Z700: 5台 → FW烧录 / 上下料 / 老化箱插拔",
    "  TARS: 5台同步部署",
    "  验收标准: 插拔成功率 ≥99%，零损伤",
    "",
    "Phase 2  (2026.12 — 2027.04)  全线部署 + 产线联调",
    "  Z700: 热海柜操作 + ATS自动测试 + 全线串联",
    "  TARS: 补充优化",
    "  验收标准: 全工站节拍达标，力控闭环验证",
    "",
    "Phase 3  (2027.05+)  规模复制 + 行业扩展",
    "  跨行业复制: 3C电子 / 汽车电子 / 半导体封测",
    "  TARS: 评估续约",
    "  验收标准: ROI验证，客户复购",
])

# ═══ Slide 9: Business ═══
add_stats_slide("商业价值", [
    ("3-5×", "", "人工效率提升"),
    ("24/7", "", "连续不间断运行"),
    ("<12月", "", "投资回收期"),
    ("5-8", "工站/台", "单台覆盖能力"),
    ("15-25", "等效人", "年化人力节省"),
    ("1000", "万元", "一期项目预算"),
])

out = '/www/wwwroot/datadrive.world/api/z700-analysis.pptx'
prs.save(out)
print(f"OK: {out}")
