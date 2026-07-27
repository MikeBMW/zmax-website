#!/usr/bin/env python3
"""Generate DDS-driven PPTX from database"""
import sqlite3, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

db = sqlite3.connect('/www/wwwroot/datadrive.world/dds.db')
cur = db.cursor()

# Read all data
co = {r[0]:r[1] for r in cur.execute('SELECT key,value FROM company').fetchall()}
kpis = {}
for r in cur.execute('SELECT * FROM kpi').fetchall():
    kpis[r[0]] = {'value':r[1],'unit':r[2],'label':r[3],'icon':r[4]}
robots = []
for r in cur.execute('SELECT * FROM robots').fetchall():
    robots.append({'id':r[0],'name':r[1],'level':r[2],'level_label':r[3],'desc':r[4],'icon':r[5],'page':r[6],'color':r[7]})
zones = []
for r in cur.execute('SELECT * FROM factory_zones').fetchall():
    zones.append({'id':r[0],'name':r[1],'color':r[2],'page':r[3],'stations':r[5]})
roadmap = []
for r in cur.execute('SELECT * FROM roadmap').fetchall():
    roadmap.append({'version':r[0],'timeline':r[1],'name':r[2],'desc':r[3],'color':r[4]})
skills_total = cur.execute('SELECT COUNT(*) FROM atomic_skills').fetchone()[0]
skills_cats = cur.execute('SELECT category,COUNT(*) FROM atomic_skills GROUP BY category ORDER BY COUNT(*) DESC').fetchall()
cur.close(); db.close()

C = RGBColor(0x00, 0xAA, 0x88)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bgw(s): s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
def tbar(s):
    sh = s.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    sh.fill.solid(); sh.fill.fore_color.rgb = C; sh.line.fill.background()
def title(s, t, y=0.3):
    tx = s.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12), Inches(0.7))
    p = tx.text_frame.paragraphs[0]; p.text = t; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = DARK
    ul = s.shapes.add_shape(1, Inches(0.6), Inches(y+0.65), Inches(2.5), Inches(0.04))
    ul.fill.solid(); ul.fill.fore_color.rgb = C; ul.line.fill.background()
def bullets(s, lines, y=1.3, sz=15):
    tx = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(5.5))
    tf = tx.text_frame; tf.word_wrap = True
    for i,ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = GRAY; p.space_after = Pt(5)
def card(s, x, y, w, h, ti, bo):
    cd = s.shapes.add_shape(1, x, y, w, h)
    cd.fill.solid(); cd.fill.fore_color.rgb = LIGHT; cd.line.color.rgb = RGBColor(0xDD,0xEE,0xEE)
    tx = s.shapes.add_textbox(x+Inches(0.15), y+Inches(0.1), w-Inches(0.3), Inches(0.35))
    tx.text_frame.paragraphs[0].text = ti; tx.text_frame.paragraphs[0].font.size = Pt(13); tx.text_frame.paragraphs[0].font.bold = True; tx.text_frame.paragraphs[0].font.color.rgb = DARK
    tx2 = s.shapes.add_textbox(x+Inches(0.15), y+Inches(0.5), w-Inches(0.3), h-Inches(0.6))
    tf2 = tx2.text_frame; tf2.word_wrap = True
    for i,ln in enumerate(bo.split('\n')):
        p = tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
        p.text = ln; p.font.size = Pt(9); p.font.color.rgb = GRAY; p.space_after = Pt(2)

# Slide 1: Cover
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
tx = s.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(4.5))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = co.get('product','Z-MAX')+' 光模块精密制造'; p.font.size = Pt(22); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Z700 轮式双臂机器人"; p2.font.size = Pt(42); p2.font.bold = True; p2.font.color.rgb = C; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
p3 = tf.add_paragraph(); p3.text = co.get('product_tag','具身智能精密制造'); p3.font.size = Pt(24); p3.font.color.rgb = DARK; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(8)
p4 = tf.add_paragraph(); p4.text = "立 项 申 请 书"; p4.font.size = Pt(36); p4.font.bold = True; p4.font.color.rgb = C; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(20)
p5 = tf.add_paragraph(); p5.text = co.get('name','')+' · '+co.get('domain','')+' · '+co.get('year','2026'); p5.font.size = Pt(14); p5.font.color.rgb = GRAY; p5.alignment = PP_ALIGN.CENTER; p5.space_before = Pt(30)
bot = s.shapes.add_shape(1, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08))
bot.fill.solid(); bot.fill.fore_color.rgb = C; bot.line.fill.background()

# Slide 2: KPI
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
title(s, "核心性能指标")
ki = list(kpis.values())
for i,k in enumerate(ki[:6]):
    x = Inches(0.5 + i*2.1); y = Inches(1.5)
    cd = s.shapes.add_shape(1, x, y, Inches(1.8), Inches(2.0))
    cd.fill.solid(); cd.fill.fore_color.rgb = LIGHT; cd.line.color.rgb = RGBColor(0xDD,0xEE,0xEE)
    tv = s.shapes.add_textbox(x, y+Inches(0.3), Inches(1.8), Inches(0.7))
    tv.text_frame.paragraphs[0].text = k['value']; tv.text_frame.paragraphs[0].font.size = Pt(28); tv.text_frame.paragraphs[0].font.bold = True; tv.text_frame.paragraphs[0].font.color.rgb = C; tv.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tu = s.shapes.add_textbox(x, y+Inches(0.95), Inches(1.8), Inches(0.35))
    tu.text_frame.paragraphs[0].text = k['unit']; tu.text_frame.paragraphs[0].font.size = Pt(11); tu.text_frame.paragraphs[0].font.color.rgb = GRAY; tu.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tl = s.shapes.add_textbox(x, y+Inches(1.3), Inches(1.8), Inches(0.45))
    tl.text_frame.paragraphs[0].text = k['label']; tl.text_frame.paragraphs[0].font.size = Pt(9); tl.text_frame.paragraphs[0].font.color.rgb = GRAY; tl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
bullets(s, ["Z700 轮式双臂精密操作机器人 — "+co.get('product_tag',''),"一期10台(5自研+5它石) · 白盒交付 · 覆盖FW烧录/上下料/老化箱/热海柜/ATS五大场景"], y=Inches(4.0))

# Slide 3: Products
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
title(s, "产品与系统架构")
for i,r in enumerate(robots[:4]):
    col=i%2; row=i//2
    card(s, Inches(0.5+col*6.3), Inches(1.3+row*2.8), Inches(5.8), Inches(2.4),
        r['name']+' ['+r['level_label']+']', r['desc'])

# Slide 4: Factory
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
title(s, "目标工厂")
ts = sum(z['stations'] for z in zones)
stats = [(str(len(zones)),'区域','产线分区'),(str(ts),'工位','全流程覆盖'),(str(skills_total),'技能','原子技能库'),('10','台','一期部署')]
for i,(v,u,l) in enumerate(stats):
    x=Inches(0.8+i*3.1); y=Inches(1.5)
    cd=s.shapes.add_shape(1,x,y,Inches(2.6),Inches(1.5))
    cd.fill.solid();cd.fill.fore_color.rgb=LIGHT;cd.line.color.rgb=RGBColor(0xDD,0xEE,0xEE)
    tv=s.shapes.add_textbox(x,y+Inches(0.2),Inches(2.6),Inches(0.6))
    tv.text_frame.paragraphs[0].text=v;tv.text_frame.paragraphs[0].font.size=Pt(24);tv.text_frame.paragraphs[0].font.bold=True;tv.text_frame.paragraphs[0].font.color.rgb=C;tv.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
    tu=s.shapes.add_textbox(x,y+Inches(0.8),Inches(2.6),Inches(0.3))
    tu.text_frame.paragraphs[0].text=u;tu.text_frame.paragraphs[0].font.size=Pt(10);tu.text_frame.paragraphs[0].font.color.rgb=GRAY;tu.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
    tl=s.shapes.add_textbox(x,y+Inches(1.05),Inches(2.6),Inches(0.3))
    tl.text_frame.paragraphs[0].text=l;tl.text_frame.paragraphs[0].font.size=Pt(8);tl.text_frame.paragraphs[0].font.color.rgb=GRAY;tl.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
zlines=[z['name']+' — '+str(z['stations'])+'工位' for z in zones]
bullets(s, zlines, y=Inches(3.5), sz=13)

# Slide 5: Skills
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
title(s, "原子技能覆盖（"+str(skills_total)+"项）")
for i,(cat,cnt) in enumerate(skills_cats[:8]):
    col=i%4;row=i//4
    x=Inches(0.5+col*3.15);y=Inches(1.5+row*2.5)
    cd=s.shapes.add_shape(1,x,y,Inches(2.8),Inches(2.0))
    cd.fill.solid();cd.fill.fore_color.rgb=LIGHT;cd.line.color.rgb=RGBColor(0xDD,0xEE,0xEE)
    tv=s.shapes.add_textbox(x,y+Inches(0.4),Inches(2.8),Inches(0.7))
    tv.text_frame.paragraphs[0].text=str(cnt);tv.text_frame.paragraphs[0].font.size=Pt(32);tv.text_frame.paragraphs[0].font.bold=True;tv.text_frame.paragraphs[0].font.color.rgb=C;tv.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
    tl=s.shapes.add_textbox(x,y+Inches(1.1),Inches(2.8),Inches(0.5))
    tl.text_frame.paragraphs[0].text=cat;tl.text_frame.paragraphs[0].font.size=Pt(11);tl.text_frame.paragraphs[0].font.color.rgb=GRAY;tl.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER

# Slide 6: Roadmap
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
title(s, "开发路线图")
for i,r in enumerate(roadmap):
    card(s, Inches(0.5), Inches(1.3+i*2.0), Inches(3.5), Inches(1.6), r['version']+': '+r['name'], '时间: '+r['timeline'])
    tx=s.shapes.add_textbox(Inches(4.3), Inches(1.4+i*2.0), Inches(8.5), Inches(1.5))
    tx.text_frame.word_wrap=True
    tx.text_frame.paragraphs[0].text=r['desc'];tx.text_frame.paragraphs[0].font.size=Pt(12);tx.text_frame.paragraphs[0].font.color.rgb=GRAY

# Slide 7: Value
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
title(s, "项目核心价值")
tx=s.shapes.add_textbox(Inches(1.5),Inches(2.2),Inches(10),Inches(3))
tf=tx.text_frame;tf.word_wrap=True
p=tf.paragraphs[0];p.text=co.get('product_tag','');p.font.size=Pt(28);p.font.bold=True;p.font.color.rgb=C;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text=str(skills_total)+"项原子技能 × 光模块全产线覆盖 × 白盒交付";p2.font.size=Pt(16);p2.font.color.rgb=GRAY;p2.alignment=PP_ALIGN.CENTER;p2.space_before=Pt(16)

# Slide 8: Thank You
s = prs.slides.add_slide(prs.slide_layouts[6]); bgw(s); tbar(s)
tx=s.shapes.add_textbox(Inches(1.5),Inches(2.5),Inches(10),Inches(2.5))
tf=tx.text_frame;tf.word_wrap=True
p=tf.paragraphs[0];p.text="Z-MAX · 具身智能精密制造";p.font.size=Pt(36);p.font.bold=True;p.font.color.rgb=C;p.alignment=PP_ALIGN.CENTER
p2=tf.add_paragraph();p2.text=co.get('name','')+' | '+co.get('domain','');p2.font.size=Pt(18);p2.font.color.rgb=GRAY;p2.alignment=PP_ALIGN.CENTER;p2.space_before=Pt(16)

out = '/www/wwwroot/datadrive.world/Z700-立项申请书.pptx'
prs.save(out)
print(f"OK: {out} ({len(prs.slides)} slides)")
