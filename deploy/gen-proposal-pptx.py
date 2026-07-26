#!/usr/bin/env python3
"""Generate Z700 Project Proposal PPTX - 12 slides"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
C = RGBColor(0x00, 0xAA, 0x88)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x22, 0x66, 0xCC)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)

def bg_white(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE

def top_bar(slide):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = C; s.line.fill.background()

def add_title(slide, text, y=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(y or 0.3), Inches(12), Inches(0.7))
    tf = tx.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = DARK
    ul = slide.shapes.add_shape(1, Inches(0.6), Inches((y or 0.3)+0.65), Inches(2.5), Inches(0.04))
    ul.fill.solid(); ul.fill.fore_color.rgb = C; ul.line.fill.background()

def add_bullets(slide, bullets, y=None, size=None):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(y or 1.3), Inches(11.5), Inches(5.5))
    tf = tx.text_frame; tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size or 15); p.font.color.rgb = GRAY
        p.space_after = Pt(5)

def add_card(slide, x, y, w, h, title, body, color=None):
    card = slide.shapes.add_shape(1, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(0xDD, 0xEE, 0xEE)
    tx = slide.shapes.add_textbox(x+Inches(0.2), y+Inches(0.15), w-Inches(0.4), Inches(0.4))
    p = tx.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = DARK
    tx2 = slide.shapes.add_textbox(x+Inches(0.2), y+Inches(0.55), w-Inches(0.4), h-Inches(0.7))
    tf2 = tx2.text_frame; tf2.word_wrap = True
    for i, line in enumerate(body.split('\n')):
        pt = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        pt.text = line; pt.font.size = Pt(10); pt.font.color.rgb = GRAY; pt.space_after = Pt(3)

def add_kpi_card(slide, x, y, val, unit, label):
    card = slide.shapes.add_shape(1, x, y, Inches(1.8), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(0xDD, 0xEE, 0xEE)
    txv = slide.shapes.add_textbox(x, y+Inches(0.3), Inches(1.8), Inches(0.7))
    pv = txv.text_frame.paragraphs[0]; pv.text = val; pv.font.size = Pt(30); pv.font.bold = True; pv.font.color.rgb = C; pv.alignment = PP_ALIGN.CENTER
    txu = slide.shapes.add_textbox(x, y+Inches(0.95), Inches(1.8), Inches(0.35))
    pu = txu.text_frame.paragraphs[0]; pu.text = unit; pu.font.size = Pt(11); pu.font.color.rgb = GRAY; pu.alignment = PP_ALIGN.CENTER
    txl = slide.shapes.add_textbox(x, y+Inches(1.3), Inches(1.8), Inches(0.45))
    pl = txl.text_frame.paragraphs[0]; pl.text = label; pl.font.size = Pt(9); pl.font.color.rgb = GRAY; pl.alignment = PP_ALIGN.CENTER

# ═══ Slide 1: Cover ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s)
top_bar(s)
tx = s.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(4.5))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Z-MAX 光模块精密制造"; p.font.size = Pt(22); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Z700 轮式双臂机器人"; p2.font.size = Pt(42); p2.font.bold = True; p2.font.color.rgb = C; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
p3 = tf.add_paragraph(); p3.text = "具身智能精密操作系统"; p3.font.size = Pt(24); p3.font.color.rgb = DARK; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(8)
p4 = tf.add_paragraph(); p4.text = "立 项 申 请 书"; p4.font.size = Pt(36); p4.font.bold = True; p4.font.color.rgb = C; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(20)
p5 = tf.add_paragraph(); p5.text = "智蜂创元 (苏州) × 它石智航 TARS (上海) · 2026"; p5.font.size = Pt(14); p5.font.color.rgb = GRAY; p5.alignment = PP_ALIGN.CENTER; p5.space_before = Pt(30)
bot = s.shapes.add_shape(1, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08))
bot.fill.solid(); bot.fill.fore_color.rgb = C; bot.line.fill.background()

# ═══ Slide 2: Overview ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "一、项目概述")
add_bullets(s, [
    "项目背景",
    "光模块是AI算力基础设施的核心器件。2026年全球光模块市场规模超200亿美元。",
    "800G/1.6T高速模块产能持续扩张，传统产线依赖人工，面临精度/一致性/招工三大痛点。",
    "",
    "项目目标",
    "研发 Z700 轮式双臂精密操作机器人，实现光模块产线关键工序的机器人替代。",
    "一期部署10台（5台Z700自研 + 5台TARS它石协同）。",
    "覆盖 FW烧录、上下料、老化箱插拔、热海柜操作、ATS自动测试 五大场景。",
])
cols = 6; cw = Inches(1.8); gap = Inches(0.15)
sx = (prs.slide_width - (cw+gap)*cols + gap) / 2
kpis = [("±0.02","mm","重复定位"),(">99","%","工序良率"),(">10","kHz","力控带宽"),("<15","s","插拔节拍"),("≤0.1","N","接触力"),(">50","kg","双臂负载")]
for i,(v,u,l) in enumerate(kpis):
    add_kpi_card(s, sx+i*(cw+gap), Inches(5.0), v, u, l)

# ═══ Slide 3: Products ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "二、产品与系统架构")
add_card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.5),
    "Z700 轮式双臂 [L4旗舰]",
    "轮式底盘自主导航 · 双臂6轴×2协同\n8+传感器融合 · 力控闭环>10kHz\n覆盖老化箱插拔/产线巡检/多工位联动\n全自主执行+自适应+自恢复")
add_card(s, Inches(6.8), Inches(1.3), Inches(5.8), Inches(2.5),
    "Z700F 固定式精密插拔 [L2基线]",
    "固定工位高精度插拔\n珞石SR5-C机械臂+DH夹爪\n5工位循环：入料→扫码→刷程序→AOI检测→出料")
add_card(s, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.8), "感知层",
    "2D/3D视觉 · 六维力 · 触觉阵列\n结构光 · 激光雷达 · IMU\n15项感知原子技能")
add_card(s, Inches(4.7), Inches(4.2), Inches(3.8), Inches(2.8), "决策层",
    "VLA大模型 · 任务规划\n路径优化 · 异常检测\n自恢复 · 8项学习泛化技能")
add_card(s, Inches(8.9), Inches(4.2), Inches(3.8), Inches(2.8), "执行层",
    "双臂协同 · 力位混合控制\n柔顺装配 · 精密抓取\n15项操作+15项装配技能")

# ═══ Slide 4: Tech ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "三、核心技术能力")
techs = [
    ("视触觉感知","3微米级缺陷识别 · 6D位姿估计\n力/触觉多模态融合\nD405深度+腕部RGB"),
    ("精密操作","力控闭环>10kHz · 力位混合\n柔顺装配 · 自适应抓取\n±0.02mm重复定位"),
    ("自主导航","AMR轮式底盘 · SLAM\n动态避障 · 工站精准对接\n多工位自由调度"),
    ("智能决策","VLA大模型推理 · 任务规划\n异常检测与自恢复\n多机协同调度"),
    ("安全集成","碰撞检测 · 力控急停\n区域监控 · 安全互锁\n10项安全原子技能"),
    ("质量追溯","AOI检测 · 3D测量 · GR&R\nMES集成 · 全流程数据闭环\n10项检测原子技能"),
]
for i,(t,b) in enumerate(techs):
    col=i%3; row=i//3
    add_card(s, Inches(0.5+col*4.1), Inches(1.3+row*2.9), Inches(3.7), Inches(2.5), t, b)

# ═══ Slide 5: Factory ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "四、目标工厂：800G DR8 光模块产线")
cols2=4;cw2=Inches(2.6);gap2=Inches(0.35)
sx2=(prs.slide_width-(cw2+gap2)*cols2+gap2)/2
zk=[("4","区域","COC+OE+MOD+WH"),("129","工位","覆盖全流程"),("89","技能","原子技能库"),("10","台","一期部署")]
for i,(v,u,l) in enumerate(zk):
    add_kpi_card(s, sx2+i*(cw2+gap2), Inches(1.5), v, u, l)
add_bullets(s, [
    "COC基板区 24站 — 打标/共晶/打线/BI老化/目检/UV固化 — Z700:FW烧录/上下料/老化箱",
    "OE光引擎区 72站 — 贴片/耦合/测试/组装/BurnIn — Z700:ATS测试/热海柜  TARS:上下料/目检",
    "MOD模块区 26站 — 模块组装/测试/包装 — Z700:全线串联",
    "WH仓库区 7站 — 来料/配送/成品/出货终检 — Z700:自主配送",
], y=Inches(3.3), size=13)

# ═══ Slide 6: Cooperation ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "五、工程合作模式：智蜂 × 它石")
add_card(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.5),
    "🐝 智蜂创元 ZFCY（甲方/苏州）",
    "AI模型与智能引擎 · 系统架构\n力控算法 · VLA大模型 · 视觉系统\n项目管理 · 客户交付\n\nASPICE: 主导SYS.1-SYS.5 · 验收SYS.6")
add_card(s, Inches(6.8), Inches(1.3), Inches(5.8), Inches(3.5),
    "💎 它石智航 TARS（乙方/上海）",
    "机械臂本体 · 运动控制 · 硬件集成\n工装治具 · 产线部署 · 现场调试\n\nASPICE: 主导SWE.1-SWE.6\n联合验收SWE.5/SYS.5")
add_bullets(s, [
    "白盒交付 · ASPICE V-Model开发体系 · 双公司协同 · 职责清晰可审计",
    "一期10台(5台Z700自研 + 5台TARS它石) · 1000万预算 · Phase1到2026.11",
], y=Inches(5.2))

# ═══ Slide 7: Roadmap ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "六、开发阶段与里程碑")
phases = [
    ("Phase 1: 单工位验证","2026.07-11",
     ["首批样机10台 · 5台Z700+5台TARS",
      "FW烧录 · 精密上下料 · 老化箱插拔",
      "验收: 插拔成功率≥99% · 力控≤0.1N · 零损伤"]),
    ("Phase 2: 全线部署","2026.12-2027.04",
     ["热海柜操作 · ATS自动测试 · 多工位串联",
      "MES全流程集成 · 力控闭环验证",
      "验收: 全工站节拍达标 · 数据闭环"]),
    ("Phase 3: 规模复制","2027.05+",
     ["跨行业复制: 3C电子/汽车电子/半导体",
      "TARS评估续约 · 标准化产品包",
      "验收: ROI验证 · 客户复购"]),
]
for i,(title,time,bullets) in enumerate(phases):
    y=Inches(1.5+i*2.0)
    add_card(s, Inches(0.5), y, Inches(3.5), Inches(1.6), title, f"时间: {time}")
    tx=s.shapes.add_textbox(Inches(4.3), y+Inches(0.1), Inches(8.5), Inches(1.5))
    tf=tx.text_frame; tf.word_wrap=True
    for j,bl in enumerate(bullets):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p.text="• "+bl; p.font.size=Pt(12); p.font.color.rgb=GRAY; p.space_after=Pt(4)

# ═══ Slide 8: Investment ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "七、投资预估与回报分析")
cols3=3;cw3=Inches(3.2);gap3=Inches(0.5)
sx3=(prs.slide_width-(cw3+gap3)*cols3+gap3)/2
inv=[("1000万","一期预算","10台机器人+软硬件+集成"),("<12月","投资回收","人工替代ROI"),("3-5×","效率提升","24/7连续运行")]
for i,(v,u,l) in enumerate(inv):
    add_kpi_card(s, sx3+i*(cw3+gap3), Inches(1.3), v, u, l)
add_bullets(s, [
    "硬件采购 ~400万: 10台机器人(5Z700+5TARS)、传感器、工装治具",
    "软件开发 ~250万: VLA模型训练、89项原子技能开发",
    "系统集成 ~150万: 产线对接、MES集成、工站改造",
    "测试验证 ~100万: POC验证、GR&R、可靠性测试",
    "项目管理 ~100万: ASPICE流程、文档、培训",
], y=Inches(3.8), size=13)

# ═══ Slide 9: Team ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "八、核心团队与技术储备")
add_card(s, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.5),
    "总工/架构 · xspace",
    "系统架构 · 4060服务器\nOrin部署 · GUI开发\nGitHub主线管理")
add_card(s, Inches(4.7), Inches(1.3), Inches(3.8), Inches(2.5),
    "产品/模型 · web",
    "产品总监 · VLA模型训练\nComfyUI前端 · 项目管理\n原子技能统筹")
add_card(s, Inches(8.9), Inches(1.3), Inches(3.8), Inches(2.5),
    "硬件/产线 · 小芳",
    "硬件集成 · MAC中转\nOrin采集 · WebSocket\n产线数据管道")
add_bullets(s, [
    "技术储备: VLA大模型(SmolVLA/GR00T/ACT/LeWM四引擎) · 力控闭环>10kHz · 3微米视觉检测",
    "89项原子技能库 · 48工位3D仿真 · A*自主导航 · DDS全局数据空间 · ASPICE V-Model开发流程",
], y=Inches(4.3))

# ═══ Slide 10: Risk ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "九、风险评估与对策")
add_bullets(s, [
    "🔴 高风险 — 力控精度不达标: 插拔损伤产品 → 分阶段验证·力传感器冗余·人工复核边界",
    "🟡 中风险 — 视觉泛化不足: 新来料识别失败 → 载具模板库·机械基准兜底·持续采集样本",
    "🟡 中风险 — 供应链延迟: 硬件交付延期 → 多供应商备选·关键件预采购·TARS产能预留",
    "🟡 中风险 — 客户接口不开放: 无法对接MES/PLC → 提前获取点表·标准协议适配",
    "🔵 低风险 — 人才流失: 研发中断 → 白盒交付·文档完善·知识沉淀到原子技能库",
], size=14)

# ═══ Slide 11: Value ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s); top_bar(s)
add_title(s, "项目核心价值")
tx = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(3.5))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "视触觉混合动作模型 · 端到端感知决策控制"; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = C; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "89项原子技能 × 光模块全产线覆盖 × 白盒交付全栈能力"; p2.font.size = Pt(16); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)
p3 = tf.add_paragraph(); p3.text = "从算法到工厂落地的完整闭环"; p3.font.size = Pt(14); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(12)
p4 = tf.add_paragraph(); p4.text = "3-5×人工效率 · 24/7连续运行 · <12月投资回收"; p4.font.size = Pt(14); p4.font.color.rgb = GRAY; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(12)

# ═══ Slide 12: Thank You ═══
s = prs.slides.add_slide(prs.slide_layouts[6]); bg_white(s)
top_bar(s)
tx = s.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(3))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Z-MAX · 具身智能精密制造"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = C; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "智蜂创元 ZFCY  |  datadrive.world"; p2.font.size = Pt(18); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)
p3 = tf.add_paragraph(); p3.text = "苏州 · 上海  |  2026"; p3.font.size = Pt(14); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(12)
bot = s.shapes.add_shape(1, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08))
bot.fill.solid(); bot.fill.fore_color.rgb = C; bot.line.fill.background()

out = '/www/wwwroot/datadrive.world/proposal.pptx'
prs.save(out)
print(f"OK: {out} ({len(prs.slides)} slides)")
