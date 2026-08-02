#!/usr/bin/env python3
"""Generate executive-briefing.pptx — 3 slides"""
import shutil, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TEMPLATE = '/www/wwwroot/datadrive.world/uploads/data/智蜂具身机器人产品规划和宣传页_v1.1_0422.pptx'
OUT = '/www/wwwroot/datadrive.world/executive-briefing.pptx'

BLUE   = RGBColor(0x00, 0x66, 0xCC)
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY  = RGBColor(0xf0, 0xf0, 0xf0)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

shutil.copy(TEMPLATE, OUT)
prs = Presentation(OUT)
layout = prs.slide_layouts[1]  # content layout with background

def add_slide(title_text):
    slide = prs.slides.add_slide(layout)
    # title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title_text
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = DARK
    return slide

def add_kpi_card(slide, left, top, val, label, sub, w=Inches(2.15), h=Inches(1.3)):
    shape = slide.shapes.add_shape(1, left, top, w, h)  # rectangle
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLUE; shape.line.width = Pt(1)
    shape.shadow.inherit = False
    # value
    txBox = slide.shapes.add_textbox(left+Inches(0.1), top+Inches(0.1), w-Inches(0.2), Inches(0.55))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = val; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER
    # label
    txBox2 = slide.shapes.add_textbox(left+Inches(0.1), top+Inches(0.6), w-Inches(0.2), Inches(0.3))
    tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
    p2.text = label; p2.font.size = Pt(12); p2.font.color.rgb = DARK
    p2.alignment = PP_ALIGN.CENTER
    # sub
    txBox3 = slide.shapes.add_textbox(left+Inches(0.1), top+Inches(0.9), w-Inches(0.2), Inches(0.35))
    tf3 = txBox3.text_frame; tf3.word_wrap = True; p3 = tf3.paragraphs[0]
    p3.text = sub; p3.font.size = Pt(9); p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

def add_table(slide, left, top, rows, col_widths, w=Inches(8.8)):
    n_rows = len(rows); n_cols = len(col_widths)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, w, Inches(0.35 * n_rows))
    tbl = tbl_shape.table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, (text, bold, color, size) in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.text = text; p.font.size = Pt(size or 11)
            p.font.bold = bold; p.font.color.rgb = color or DARK
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
                p.font.color.rgb = WHITE
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LGRAY

# ═══ Slide 1: 立项逻辑 ═══
s1 = add_slide('为什么必须先"现场演练"才能立项？')

# Left column
txBox = s1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.0), Inches(0.3))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = 'ROI需要真实数据'; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = BLUE

left_rows = [
    [('问题', True, WHITE, 10), ('说明', True, WHITE, 10)],
    [('ROI需要真实数据', True, DARK, 10), ('投资回报率不是靠理论推算，需要现场的真实节拍、成功率、人力替代比', False, GRAY, 9)],
    [('模型需要真实数据', True, DARK, 10), ('具身大模型的泛化能力依赖真机数据采集——仿真数据无法替代产线真实工况', False, GRAY, 9)],
    [('行业共识', True, DARK, 10), ('2026年6月工信部、国资委启动"实景实训专项行动"', False, GRAY, 9)],
    [('破局路径', True, DARK, 10), ('用机器人去现场演练→采集真机数据→验证技术指标→算清ROI→正式立项', False, GRAY, 9)],
]
add_table(s1, Inches(0.5), Inches(1.6), left_rows, [Inches(1.6), Inches(2.6)], Inches(4.2))

# Right column: Davinci comparison
txBox2 = s1.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(0.3))
tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
p2.text = '进工厂干活，真的非得是人形吗？'; p2.font.size = Pt(16); p2.font.bold = True; p2.font.color.rgb = ORANGE

right_rows = [
    [('对比维度', True, WHITE, 10), ('达芬奇(手术)', True, WHITE, 10), ('轮式双臂机器人', True, WHITE, 10), ('工业专用', True, WHITE, 10)],
    [('形态', True, DARK, 10), ('非人形(多臂)', False, GRAY, 9), ('轮式底盘+双臂', False, GRAY, 9), ('专用形态', False, BLUE, 9)],
    [('精度', True, DARK, 10), ('0.1mm级', False, GRAY, 9), ('~0.05-0.5mm', False, GRAY, 9), ('±0.05mm', False, BLUE, 9)],
    [('售价', True, DARK, 10), ('~3000万', False, GRAY, 9), ('15-60万', False, GRAY, 9), ('<30万', False, BLUE, 9)],
    [('核心壁垒', True, DARK, 10), ('1500万例数据', False, GRAY, 9), ('移动操作协同', False, GRAY, 9), ('场景数据+工艺', False, BLUE, 9)],
]
add_table(s1, Inches(5.2), Inches(1.6), right_rows, [Inches(1.0), Inches(1.05), Inches(1.0), Inches(1.15)], Inches(4.2))

# Bottom quote
txBox3 = s1.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9.0), Inches(0.7))
tf3 = txBox3.text_frame; tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = '达芬奇不做成人形，依然统治手术机器人市场60%份额。它的价值在于精度、灵活性、数据积累，而不是"长得像人"。'
p3.font.size = Pt(12); p3.font.italic = True; p3.font.color.rgb = GRAY

# ═══ Slide 2: 产品之道 ═══
s2 = add_slide('Z-MAX 的产品之道')
txBox4 = s2.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9.0), Inches(0.3))
tf4 = txBox4.text_frame; p4 = tf4.paragraphs[0]
p4.text = '以终为始，从场景倒推形态——光模块插拔需要什么，我们就做什么'
p4.font.size = Pt(14); p4.font.italic = True; p4.font.color.rgb = GRAY

kpis = [
    ('99%', '插拔成功率', '模型训练+真机数据闭环'),
    ('<15s', '单颗节拍', '力控算法+宏微复合控制'),
    ('±0.05mm', '定位精度', '视觉+力觉多模态融合'),
    ('±0.1N', '力控精度', '六维力传感器+阻抗控制'),
    ('≤0.5h', '换型时间', '软件定义，无需二次编程'),
    ('24h', '连续运行', '硬件可靠性+异常自恢复'),
]
for i, (v, l, s) in enumerate(kpis):
    col = i % 3; row = i // 3
    add_kpi_card(s2, Inches(0.5 + col * 3.1), Inches(1.55 + row * 1.55), v, l, s)

# Strategy cards as text boxes
cards_data = [
    ('🧠 智能驱动', 'VLA-Touch端到端感知-决策-控制大模型。真机数据闭环持续进化，每10万次插拔精度提升0.5%。'),
    ('🔧 硬件本体', '7轴力控机械臂+D405深度相机+六维力传感器+专用末端执行器。不做通用人形，做最优形态。'),
    ('📊 数据闭环', '采集→中转→训练→验证→部署五阶段全链路。每颗插拔数据自动回流，模型周级迭代。'),
    ('🏭 落地路径', '800G DR8光模块产线128站→先验证关键工位→逐站推广。3个月出数据，6个月算清ROI。'),
]
for i, (title, desc) in enumerate(cards_data):
    col = i % 2; row = i // 2
    x = Inches(0.5 + col * 4.8); y = Inches(4.75 + row * 0.9); w = Inches(4.4)
    shape = s2.shapes.add_shape(1, x, y, w, Inches(0.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd); shape.line.width = Pt(0.5)
    txB = s2.shapes.add_textbox(x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.35))
    t = txB.text_frame; pp = t.paragraphs[0]
    pp.text = title; pp.font.size = Pt(12); pp.font.bold = True; pp.font.color.rgb = BLUE
    txB2 = s2.shapes.add_textbox(x + Inches(0.15), y + Inches(0.35), w - Inches(0.3), Inches(0.4))
    t2 = txB2.text_frame; t2.word_wrap = True; pp2 = t2.paragraphs[0]
    pp2.text = desc; pp2.font.size = Pt(9); pp2.font.color.rgb = GRAY

# ═══ Slide 3: 市场预判与投入 ═══
s3 = add_slide('现在处于哪个阶段？马上会发生什么？')
txBox5 = s3.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9.0), Inches(0.3))
tf5 = txBox5.text_frame; p5 = tf5.paragraphs[0]
p5.text = '从"讲故事"到"交付"的分水岭——不进则退的窗口期'
p5.font.size = Pt(14); p5.font.italic = True; p5.font.color.rgb = GRAY

# Left: Market data
market_rows = [
    [('维度', True, WHITE, 10), ('现状', True, WHITE, 10)],
    [('行业共识', True, DARK, 10), ('2025年量产元年，2026-2027年从"讲故事"转向"交付时代"', False, GRAY, 9)],
    [('政策驱动', True, DARK, 10), ('具身智能连续两年写入政府工作报告；实景实训专项行动已启动', False, GRAY, 9)],
    [('市场规模', True, DARK, 10), ('2025年9150亿元，2026年预计突破1万亿元', False, GRAY, 9)],
    [('资本热度', True, DARK, 10), ('2026上半年融资935亿元，较2025年同期提升5倍', False, GRAY, 9)],
    [('行业痛点', True, DARK, 10), ('"缺数据是共识"——真机数据是最大瓶颈', False, GRAY, 9)],
]
add_table(s3, Inches(0.5), Inches(1.5), market_rows, [Inches(1.2), Inches(3.2)], Inches(4.4))

# Right: Investment
invest_rows = [
    [('性能指标', True, WHITE, 10), ('目标', True, WHITE, 9), ('投入方向', True, WHITE, 10), ('预估', True, WHITE, 10)],
    [('插拔成功率', True, DARK, 10), ('95→99%', True, BLUE, 9), ('真机数据采集10万次+模型训练+闭环迭代', False, GRAY, 8), ('180万', True, ORANGE, 10)],
    [('单颗节拍', True, DARK, 10), ('30→15s', True, BLUE, 9), ('力控算法优化+宏微复合控制+高速相机标定', False, GRAY, 8), ('120万', True, ORANGE, 10)],
    [('定位精度', True, DARK, 10), ('±0.05mm', True, BLUE, 9), ('D405深度相机+视觉标定+多模态融合', False, GRAY, 8), ('90万', True, ORANGE, 10)],
    [('力控精度', True, DARK, 10), ('±0.1N', True, BLUE, 9), ('六维力传感器×2+阻抗控制+标定工装', False, GRAY, 8), ('110万', True, ORANGE, 10)],
    [('换型时间', True, DARK, 10), ('≤0.5h', True, BLUE, 9), ('软件定义平台+视觉模板匹配+自动标定', False, GRAY, 8), ('70万', True, ORANGE, 10)],
    [('连续运行', True, DARK, 10), ('24h', True, BLUE, 9), ('硬件可靠性测试+异常检测+自恢复', False, GRAY, 8), ('60万', True, ORANGE, 10)],
    [('现场部署', True, DARK, 10), ('3站点', True, BLUE, 9), ('Z700本体×3+现场集成+产线对接+驻场', False, GRAY, 8), ('360万', True, ORANGE, 10)],
    [('合计', True, WHITE, 11), ('12个月', True, WHITE, 10), ('Phase1(6个月)+Phase2(6个月)', True, WHITE, 9), ('990万', True, WHITE, 13)],
]
add_table(s3, Inches(5.0), Inches(1.5), invest_rows,
          [Inches(1.15), Inches(0.85), Inches(2.0), Inches(0.7)], Inches(4.7))
# Color total row
tbl = s3.shapes[-1].table
for ci in range(4):
    cell = tbl.cell(8, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = BLUE

# Bottom summary
txBox6 = s3.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9.0), Inches(0.8))
tf6 = txBox6.text_frame; tf6.word_wrap = True
p6 = tf6.paragraphs[0]
p6.text = '一句话总结：不进工厂→拿不到真机数据→训不出好模型→交不出货→被淘汰。尽快推动Z700进工厂演练，3个月出数据，6个月算清ROI，12个月全产线覆盖。'
p6.font.size = Pt(11); p6.font.bold = True; p6.font.color.rgb = ORANGE

prs.save(OUT)
print(f'done: {OUT}')
